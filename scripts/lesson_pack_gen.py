#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
lesson-pack 档B 生成器：把单源内容(JSON)落成成套真实 Office 文件。
产物：
  教学设计.docx   教学课件.pptx   学习任务单.docx   课堂分层练习.docx   课后作业.docx
数据驱动、任意学科/课题通用；PPT 为数据驱动通用课件(页数随环节动态)。

用法（推荐，外部单源 JSON）：
  python lesson_pack_gen.py --out <输出目录> --data <内容.json> [--theme 学科|'#hex']
  内容 JSON 结构见 ../templates/lesson-content.schema.md，
  可直接跑示例见 ../templates/示例_语文_七年级_咏雪_40分钟.json。
不传 --data 时回退到内置数学《一元一次方程》样例（脚本底部 LESSON 等四结构）。
"""
import os, sys, argparse, json, colorsys
from docx import Document
from docx.shared import Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PPt
    from pptx.dml.color import RGBColor as PRGB
    from pptx.enum.text import PP_ALIGN
    PPT_OK = True
except Exception:
    PPT_OK = False

RED = RGBColor(0xB3, 0x26, 0x1E)   # 红笔批注 #b3261e
BLUE = RGBColor(0x1A, 0x4F, 0x8B)  # 蓝栏 #1a4f8b
INK = RGBColor(0x2A, 0x28, 0x26)
GRAY = RGBColor(0x55, 0x55, 0x55)

# ============================================================
# 主题系统：让整套产物(教案/课件/卷子)的「主色」可自定义。
# 内置若干学科预设，也支持任意十六进制主色。
# 每个主题：head=主标题/标题栏色  acc=强调  pale=主色淡底
#            pale_r=警示淡底  red=红笔批注
# ============================================================
THEMES = {
    "数学":   {"head":"0F3B6B", "acc":"1A4F8B", "pale":"EEF3FA", "pale_r":"FCEEED", "red":"B3261E", "acc2":"0E7490"},
    "语文":   {"head":"7A1F2B", "acc":"9A3340", "pale":"FBF0F1", "pale_r":"FCEEED", "red":"B3261E", "acc2":"B4441E"},
    "英语":   {"head":"1F4E79", "acc":"2E6FA3", "pale":"EAF2F9", "pale_r":"FCEEED", "red":"B3261E", "acc2":"1E8449"},
    "物理":   {"head":"0E5A6D", "acc":"11707F", "pale":"E9F3F5", "pale_r":"FCEEED", "red":"B3261E", "acc2":"3A6E1A"},
    "化学":   {"head":"2F4858", "acc":"466E88", "pale":"EEF1F4", "pale_r":"FCEEED", "red":"B3261E", "acc2":"8B5E00"},
    "生物":   {"head":"1F5F3B", "acc":"2E7D4F", "pale":"EAF4EF", "pale_r":"FCEEED", "red":"B3261E", "acc2":"0E7490"},
    "历史":   {"head":"6D3B1F", "acc":"8A5233", "pale":"F7F0E9", "pale_r":"FCEEED", "red":"B3261E", "acc2":"B4441E"},
    "地理":   {"head":"3C5A78", "acc":"51718F", "pale":"EEF3F7", "pale_r":"FCEEED", "red":"B3261E", "acc2":"6E4A1E"},
    "通用":   {"head":"1F3B5C", "acc":"2E568A", "pale":"EDF1F6", "pale_r":"FCEEED", "red":"B3261E", "acc2":"3A6E1A"},
}

def _px(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


def _normalize_hex(theme):
    """把 --theme 的非预设取值规范成 6 位大写 hex；非法则返回 None。
    自定义主色真正参与调色板派生后，非法输入（如 --theme XX未知色）会在 int() 处崩溃，
    这里提前拦下并给出可用取值提示。"""
    raw = theme.lstrip('#').strip()
    if len(raw) == 3:                       # 支持 #C39 简写 → CC3399
        raw = ''.join(c * 2 for c in raw)
    if len(raw) != 6:
        return None
    try:
        int(raw, 16)
    except ValueError:
        return None
    return raw.upper()


def _mix_white(hexcolor, white_ratio=0.88):
    """把主色朝白色混合，得到淡底（如 pale）。white_ratio=0.88 表示白占 88%。"""
    h = hexcolor.lstrip('#')
    rgb = [int(h[i:i+2], 16) for i in (0, 2, 4)]
    return ''.join(f"{round(c * (1 - white_ratio) + 255 * white_ratio):02X}" for c in rgb)


def _rotate_hue(hexcolor, deg):
    """色相旋转得到辅色（acc2），保持明度/饱和度不变。"""
    h = hexcolor.lstrip('#')
    r, g, b = [int(h[i:i+2], 16) / 255 for i in (0, 2, 4)]
    hh, ll, ss = colorsys.rgb_to_hls(r, g, b)
    hh = (hh + deg / 360.0) % 1.0
    r, g, b = colorsys.hls_to_rgb(hh, ll, ss)
    return ''.join(f"{round(c * 255):02X}" for c in (r, g, b))

def _ppx(h):
    h = h.lstrip('#')
    return PRGB(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# 当前生效主题（运行时由 --theme 设定）
CUR = None
def apply_theme(theme):
    """theme: 预设名(数学/语文…) 或 '#hex主色'。返回主题 dict(docx用RGBColor，ppt用PRGB)。
    若给自定义主色，则从主色自动推导一套淡色调色板。"""
    global CUR
    t = THEMES.get(theme)
    if t is None:
        # 当作自定义主色：真正派生一套淡色板（旧实现 pale/acc2 是写死的，
        # 与 SKILL.md 里"自动派生淡色板"的说法不符——现在名副其实）
        base = _normalize_hex(theme)
        if base is None:
            raise SystemExit(
                f"✗ --theme 取值非法：{theme!r}。"
                f"可用预设：{'/'.join(THEMES)}；或 6 位十六进制主色，如 '#C0392B'（也支持 #C39 简写）。")
        t = {"head": base, "acc": base,
             "pale":   _mix_white(base, 0.88),   # 主色 12% + 白 88% → 极淡底
             "pale_r": "FCEEED",                 # 警示淡底为约定色，不随主题
             "red":    "B3261E",                 # 红笔批注为约定色，不随主题
             "acc2":   _rotate_hue(base, 150)}   # 色相旋转得到辅色
    d = {k: _px(v) for k,v in t.items()}      # docx RGBColor
    # ppt PRGB（仅当 python-pptx 可用；否则该侧为空，不影响 docx 生成）
    p = {}
    if PPT_OK:
        p = {k: _ppx(v) for k,v in t.items()}
    # docx 附加中性色（红笔为约定警示色，各主题统一固定）
    d['red']  = RED        # 红笔批注（约定色，不随主题）
    d['ink']  = INK
    d['gray'] = GRAY
    # ppt 附加中性色
    if PPT_OK:
        p['red']   = PRGB(0xB3,0x26,0x1E)   # 红笔警示（约定色）
        p['ink']   = PRGB(0x24,0x24,0x24)
        p['sub']   = PRGB(0x55,0x55,0x55)
        p['white'] = PRGB(0xFF,0xFF,0xFF)
        p['tag']   = PRGB(0xC9,0xD8,0xEA)
    CUR = {"docx":d, "ppt":p, "name":theme}
    return CUR

# docx 侧"蓝栏"栏目色——随主题；未指定时退化为数学深蓝(与旧默认一致)
def _acc():
    if CUR is not None:
        return CUR['docx']['acc']
    return RGBColor(0x1A, 0x4F, 0x8B)

def set_cn(run, font='楷体'):
    """设置中文字体(楷体优先，贴近纸质教案质感)"""
    run.font.name = font
    r = run._element
    rPr = r.get_or_add_rPr()
    rFonts = rPr.find(qn('w:rFonts'))
    if rFonts is None:
        rFonts = OxmlElement('w:rFonts'); rPr.append(rFonts)
    rFonts.set(qn('w:eastAsia'), font)

def para(doc, text='', size=12, bold=False, color=INK, align=None,
         cn='楷体', space_after=4, line=None, indent=None):
    p = doc.add_paragraph()
    if align: p.alignment = align
    if indent is not None: p.paragraph_format.left_indent = Cm(indent)
    p.paragraph_format.space_after = Pt(space_after)
    if line: p.paragraph_format.line_spacing = line
    run = p.add_run(text)
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    set_cn(run, cn)
    return p

def rich_para(doc, segments, size=12, align=None, space_after=4, line=1.5):
    """segments: list of (text, bold, color)"""
    p = doc.add_paragraph()
    if align: p.alignment = align
    p.paragraph_format.space_after = Pt(space_after)
    if line: p.paragraph_format.line_spacing = line
    for txt, bold, col in segments:
        run = p.add_run(txt)
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = col
        set_cn(run)
    return p

def title(doc, text, size=20, space_after=6):
    para(doc, text, size=size, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, space_after=space_after, cn='宋体')

def heading(doc, text, size=15, color=INK, space_after=4):
    para(doc, text, size=size, bold=True, color=color, space_after=space_after, cn='黑体')

def cell_text(cell, text, bold=False, size=11.5, color=INK, cn='楷体'):
    cell.text = ''
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = Pt(2)
    run = p.add_run(text)
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    set_cn(run, cn)

def hline(cell):
    """单元格底色(浅米色)"""
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd'); shd.set(qn('w:val'),'clear'); shd.set(qn('w:fill'),'F6F2E6')
    tcPr.append(shd)

# ============================================================
# 教学内容数据（单源）。此结构是档B的唯一来源，与档A内容一致。
# ============================================================
LESSON = {
    "学科": "数学", "年级": "七年级", "教材": "人教版七上",
    "章节": "第三章", "课题": "一元一次方程", "课型": "新授课",
    "课时": 45, "课时标签": "第 1 课时",
    "教学目标": [
        ("知识与技能", ["说出方程、方程的解的概念，能识别一元一次方程的三个特征；",
                       "掌握解“ax+b=c”型方程的基本步骤（移项、合并、系数化 1），正确求解。"]),
        ("过程与方法", ["通过列方程解实际问题，体会“设未知数—找等量关系—列方程”的建模过程；通过对比解法，归纳移项法则。"]),
        ("情感态度", ["感受方程解决实际问题的便利，愿意用方程表达数量关系。"]),
    ],
    "重难点": ("一元一次方程的概念识别，标准解法步骤。",
               "①从实际问题中找出等量关系并正确列式；②理解移项变号背后是等式性质。",
               "用天平演示“两边同加同减”，先直观后抽象；列式环节放慢，先让学生“说等量关系”再写方程。"),
    "学情": "学生在小学已学过简易方程（如 ax±b=c），会四则运算；多数能照步骤解，但对移项为何变号理解偏浅，遇带负号的项易错，从文字列方程的能力普遍较弱。",
    "时间总表": [
        ("课堂导入", 4), ("新知讲授", 18), ("当堂巩固练习", 15), ("课堂小结", 5), ("作业布置与收尾", 3),
    ],
    "环节": [
        {"名": "一、课堂导入", "分": 4, "rows": [
            ("教师活动", "出情境：小明的零花钱是去年的 2 倍还多 3 元，今年共 13 元，去年多少？让学生用算式和方程两种方法表示。"),
            ("学生活动", "口答结果，尝试把等量关系写成等式。"),
            ("设计意图", "从熟悉的生活问题切进，唤起简易方程旧知，引出“方程”。"),
            ("当堂检验", "抽 2–3 人说出所列等式，看是否已会找等量关系。"),
        ]},
        {"名": "二、新知讲授：概念 + 解法", "分": 18, "rows": [
            ("教师活动", "①下定义：含未知数的等式叫方程，逐条讲“一元一次”（一个未知数、次数 1、整式）；②用天平演示移项，点明“移项变号”来自等式性质；③板书标准解法：移项→合并同类项→系数化 1→检验，示范例 1 完整过程。"),
            ("学生活动", "判断一组式子哪个是方程、哪个是一元一次方程；跟随书写例 1，并说每步依据。"),
            ("设计意图", "先概念后解法，符合认知顺序；天平打通算理，避免死记。"),
            ("当堂检验", "给 2 个式子快速判断（举手/口答），确认概念真掌握再往下走。"),
        ]},
        {"名": "三、当堂巩固练习", "分": 15, "rows": [
            ("教师活动", "分层巡堂：全班先做基础 3 题，再引导做提升，最后给有余力者挑战题；巡视纠错、个别点拨。"),
            ("学生活动", "独立完成基础题，核对；有余力者继续挑战题。"),
            ("设计意图", "当堂检测 + 分层，不同程度学生都有事做、有收获。"),
            ("当堂检验", "巡堂记录正确率，收集典型错法；抽 1–2 题请学生上台板演互评。"),
        ]},
        {"名": "四、课堂小结", "分": 5, "rows": [
            ("教师活动", "带学生回顾：认识了什么？会做什么？再强调“移项要变号”这一易错点，出示本课知识结构。"),
            ("学生活动", "每人用一句话说今天学会什么；在纸条写下“今天最容易错的是____”交上。"),
            ("当堂检验", "回收小纸条，即知薄弱点集中在哪。"),
        ]},
        {"名": "五、作业布置与收尾", "分": 3, "rows": [
            ("教师活动", "分层布置作业：A 组必做、B 组选做、C 组挑战；预告下节讲去括号、去分母的方程。"),
            ("学生活动", "记清必做/选做范围。"),
            ("设计意图", "作业量控制在 20 分钟以内，不与课时争夺精力。"),
        ]},
    ],
    "板书": {
        "title": "一元一次方程",
        "concept": "概念：含未知数的等式；一元一次：一个未知数 · 次数 1 · 整式",
        "left": ["解法步骤", "移项（变号）", "合并同类项", "系数化 1", "检验"],
        "mid": ["例 1：3x−7=20", "移项 3x=20+7", "3x=27", "x=9", "检验 3×9−7=20 ✓"],
        "right": ["[留白区", "学生板演]"],
    },
    "作业说明": "预计 15–20 分钟完成；作业量不超课时一半，避免负担。",
}

# 学案/练习/作业题（卷子用，含班级姓名栏）
WORKSHEET_TASKS = [
    ("任务一 · 判断概念", "1", "下面哪些是方程？哪些是一元一次方程？（写编号）① x+2　② 3y=9　③ 2x²=8　④ 5−x=2　⑤ 1/x=3\n是方程：____　是一元一次方程：____"),
    ("任务二 · 跟着解方程", "2", "解 3x−7=20，填完整：① 移项：3x=20__7（填＋或−）② 合并：3x=__　③ 系数化1：x=__　④ 检验：3×__−7=__，正好=20。"),
    ("任务三 · 独立解（基础，必做）", "3", "解下列方程：① 2x+5=13　② 4y−3=9　③ 6−m=2"),
    ("任务四 · 挑战（选做）", "4", "小明的年龄乘 3 再减 5，正好等于 34 减 3，求小明年龄。（设 x→找等量关系→列方程）"),
    ("任务五 · 小结自评", "5", "我今天学会了：____　我还想问：____"),
]

EXERCISES = [
    ("★ 全班必做", [
        ("1", "下列是一元一次方程的是（　　）A. x²=4　B. 3x+2　C. 2y−1=5　D. 1/x=2", "C"),
        ("2", "解 3x+6=18，移项后正确的是（　　）A. 3x=18+6　B. 3x=18−6　C. 3x=−6　D. 3x=6", "B"),
        ("3", "解方程并检验：5x − 4 = 21", "5x=21+4→5x=25→x=5，检验 5×5−4=21 ✓"),
    ]),
    ("★ 能力提升（大部分同学做）", [
        ("4", "解方程：7 − 2x = 1", "−2x=1−7→−2x=−6→x=3"),
        ("5", "一个数的 3 倍比它的 2 倍大 5，求这个数（列方程求解）。", "设该数为 x，3x=2x+5→x=5"),
    ]),
    ("★ 思维挑战（学有余力选做）", [
        ("6", "自己编一道“ax+b=c”形式的实际问题让同桌列方程；并想想：解方程为什么最后都要检验？", "开放题；检验是因化归过程可能产生使原方程无意义的解（下节去分母体会）。"),
    ]),
]

HOMEWORK = {
    "A 组 · 基础必做": [
        ("1", "指出下面哪些是一元一次方程，并说明理由。", "略（考查概念三要素）"),
        ("2", "解方程：① 4x+7=23　② 2−3y=8", "① x=4　② −2y=6→y=−3"),
        ("3", "妈妈今年 36 岁，正好是小明年龄的 3 倍，小明几岁？（列方程）", "设小明 x 岁，3x=36→x=12"),
    ],
    "B 组 · 提升选做": [
        ("4", "解方程：9 − 5x = 2x − 5", "9+5=2x+5x→14=7x→x=2"),
        ("5", "一个两位数，十位与个位数字和为 9，个位是十位的 2 倍，求这个数。", "设十位 x，个位 2x，x+2x=9→x=3，数为 36"),
    ],
    "C 组 · 拓展挑战（衔接下节）": [
        ("6", "预习：解 (x+1)/2=3 时，为什么可以两边同乘 2 去掉分母？用等式性质说明并解出。", "等式性质2，两边同乘2：x+1=6→x=5"),
    ],
}

# 上面四个结构仅作「不传 --data」时的演示样例。传入外部 JSON 后 USING_SAMPLE=False，
# 缺失的键一律跳过对应产物，绝不回填内置样例（详见 load_data 说明）。
USING_SAMPLE = True

def check_time(L):
    """时间账硬校验（quality-checklist A 项）：时间总表求和 == 课时，且环节分求和也 == 课时。
    返回 (ok, 问题列表)。不抛异常——由调用方决定告警还是失败。"""
    issues = []
    total = None
    try:
        total = sum(int(m) for _, m in L.get('时间总表') or [])
    except Exception as e:
        issues.append(f"时间总表格式异常（应为 [[环节名, 分钟], …]）：{e}")
    if total is not None:
        try:
            if total != int(L.get('课时')):
                issues.append(f"时间总表合计 {total}′ ≠ 课时 {L.get('课时')}′")
        except Exception:
            issues.append(f"课时字段缺失或非数字：{L.get('课时')!r}")
    # 环节分之和（schema 硬规则第 2 条，此前只校验了时间总表，这里补齐）
    try:
        ph_sum = sum(int(p['分']) for p in L.get('环节') or [])
        if ph_sum != int(L.get('课时')):
            issues.append(f"各环节分钟合计 {ph_sum}′ ≠ 课时 {L.get('课时')}′")
    except Exception as e:
        issues.append(f"环节分钟（环节[].分）格式异常：{e}")
    return (len(issues) == 0), issues


def load_data(path):
    """从外部单源 JSON 载入内容，实现任意学科/课题通用化。
    JSON 顶层四个键与渲染逻辑解耦（结构见 templates/lesson-content.schema.md）：
      lesson           → LESSON            (dict，含学科/年级/课题/环节/板书…)
      worksheet_tasks  → WORKSHEET_TASKS   ([[任务名, 序号, 题面…]])
      exercises        → EXERCISES         ([[组名, [[题号,题面,答案]…]], …])
      homework         → HOMEWORK          ({组名: [[题号,题面,答案]…]})

    ⚠ 关键约定：一旦传入外部 JSON，**绝不用内置数学样例回填缺失的键**。
    旧实现是"只覆盖 JSON 里存在的键"，缺 worksheet_tasks 时保留内置《一元一次方程》，
    结果产出"语文《咏雪》的学案里全是方程题"的串味事故且无任何提示。
    现在的策略：缺哪个键 → 对应产物直接跳过并明确告知，宁缺勿错。
    内置样例只在「不传 --data」的演示模式下生效。
    """
    global LESSON, WORKSHEET_TASKS, EXERCISES, HOMEWORK, USING_SAMPLE
    with open(path, encoding='utf-8') as f:
        j = json.load(f)
    if 'lesson' not in j:
        raise SystemExit("✗ 内容 JSON 缺顶层 'lesson' 键（结构见 templates/lesson-content.schema.md）")
    LESSON = j['lesson']
    WORKSHEET_TASKS = j.get('worksheet_tasks', [])
    EXERCISES = j.get('exercises', [])
    HOMEWORK = j.get('homework', {})
    USING_SAMPLE = False
    return LESSON

def strip_prefix(text, prefixes):
    """去掉数据里自带的"重点：/难点：/突破："类前缀。
    渲染时脚本统一加前缀，若数据本身已带（示例 JSON 就是），不去重会输出
    "重点：重点：疏通文意…"。此处做幂等处理，两种写法都能正确渲染。"""
    t = str(text).strip()
    for p in prefixes:
        if t.startswith(p):
            return t[len(p):].strip()
    return t


def chapter_text(L):
    """章节文案：数据里常写成"第三章""第二单元第8课"，已自带"第"字。
    旧实现硬编码 f"第{章节}" → 输出"第第三章"。这里做幂等：已有"第"开头就不再加。"""
    ch = str(L.get('章节', '')).strip()
    if not ch:
        return ch
    return ch if ch.startswith('第') else f"第{ch}"


def make_design_docx(path):
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = sec.bottom_margin = Cm(2.2)
    sec.left_margin = sec.right_margin = Cm(2.5)
    L = LESSON

    title(doc, "教 学 设 计")
    para(doc, f"{L['教材']} · {chapter_text(L)} ·  {L['课题']}（{L['课时标签']} · {L['课型']}）",
         size=11, color=GRAY, align=WD_ALIGN_PARAGRAPH.CENTER, space_after=10, cn='宋体')

    # 表头信息
    t = doc.add_table(rows=2, cols=4); t.style = 'Table Grid'; t.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr = [("学科 / 年级", f"{L['学科']} · {L['年级']}"), ("课型", L['课型']),
           ("教材位置", f"{L['教材']} {chapter_text(L)}"), ("课时", f"{L['课时']} 分钟")]
    for i,(k,v) in enumerate(hdr):
        cell_text(t.cell(i//2, (i%2)*2), k, bold=True); hline(t.cell(i//2, (i%2)*2))
        cell_text(t.cell(i//2, (i%2)*2+1), v)
    doc.add_paragraph()

    # 目标
    heading(doc, "一、教学目标")
    for cat, items in L['教学目标']:
        p = doc.add_paragraph(); p.paragraph_format.space_after = Pt(2)
        run = p.add_run(f"【{cat}】"); run.font.bold=True; run.font.color.rgb=_acc(); run.font.size=Pt(12.5); set_cn(run)
        for it in items:
            rich_para(doc, [("• "+it, False, INK)], size=12)

    # 重难点（前缀由脚本统一加；数据若已带前缀则先去重，避免"重点：重点：…"）
    heading(doc, "二、教学重难点")
    rich_para(doc, [("重点：", True, INK),
                    (strip_prefix(L['重难点'][0], ("重点：", "重点:")), False, INK)])
    rich_para(doc, [("难点：", True, INK),
                    (strip_prefix(L['重难点'][1], ("难点：", "难点:")), False, INK)])
    rich_para(doc, [("突破：", True, RED),
                    (strip_prefix(L['重难点'][2], ("突破：", "突破:", "难点突破：", "难点突破策略：")), False, RED)])

    # 学情
    heading(doc, "三、学情分析")
    para(doc, L['学情'], size=12)

    # 时间
    heading(doc, "四、教学环节")
    ts = " ｜ ".join(f"{n} {m}′" for n,m in L['时间总表'])
    rich_para(doc, [("时间分配：", True, INK), (ts, False, INK),
                    ("　—— 合计 "+str(L['课时'])+"′", True, RED)])

    # 环节表格
    for ph in L['环节']:
        rich_para(doc, [(f"{ph['名']}　", True, INK), (f"{ph['分']}′", True, RED)], size=13, space_after=4)
        tb = doc.add_table(rows=len(ph['rows']), cols=2); tb.style='Table Grid'
        tb.columns[0].width = Cm(2.4); tb.columns[1].width = Cm(13.6)
        for ri,(k,v) in enumerate(ph['rows']):
            cell_text(tb.cell(ri,0), k, bold=True, size=11); hline(tb.cell(ri,0))
            cell_text(tb.cell(ri,1), v, size=11)
        doc.add_paragraph()

    # 板书
    heading(doc, "五、板书设计")
    board = doc.add_table(rows=1, cols=3); board.style='Table Grid'
    # 各分区用 .get 兜底：concept 缺失时旧实现直接 KeyError，导致整份教案生成失败
    b = L.get('板书') or {}
    left_txt = "\n".join([b.get('title') or L['课题'], b.get('concept') or ''] + list(b.get('left') or []))
    cell_text(board.cell(0,0), left_txt.strip(), size=11)
    cell_text(board.cell(0,1), "\n".join(b.get('mid') or []), size=11)
    cell_text(board.cell(0,2), "\n".join(b.get('right') or []), size=11)
    doc.add_paragraph()
    para(doc, "—— 教学设计 · 完 ——", size=10.5, color=GRAY, align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.save(path)
    return path

def make_worksheet_docx(path):
    doc = Document(); sec = doc.sections[0]
    sec.top_margin = sec.bottom_margin = Cm(2); sec.left_margin = sec.right_margin = Cm(2.5)
    title(doc, "学 习 任 务 单", size=18)
    para(doc, LESSON['课题'] + " · 课堂用", size=11, color=GRAY, align=WD_ALIGN_PARAGRAPH.CENTER)
    para(doc, "班级：＿＿＿＿　姓名：＿＿＿＿　日期：＿＿＿＿　　座号：＿＿＿＿", size=11, space_after=8)
    # 学习目标：从单源教学目标动态生成，去数学硬编码
    _goals = []
    for _cat, _its in LESSON['教学目标']:
        for _g in _its[:2]:                     # 每类最多取前2条
            if len(_goals) < 4: _goals.append(_g)
    _goal_line = "学完我能——" + "　".join("□ " + g for g in _goals)
    rich_para(doc, [("学习目标：", True, _acc()), (_goal_line, False, INK)], size=11.5)
    for hname, no, body in WORKSHEET_TASKS:
        heading(doc, hname, size=13)
        lines = body.split("\n")
        rich_para(doc, [(f"{no}. ", True, INK), (lines[0], False, INK)], size=12)
        for extra in lines[1:]:
            para(doc, extra, size=12, space_after=14)
    para(doc, "—— 学习任务单 · 发给学生当堂用 ——", size=10.5, color=GRAY, align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.save(path); return path

def exercise_minutes(L):
    """课堂练习卷的建议用时。取值优先级：
    1) lesson.练习时长（显式指定）；
    2) 名称含"练习/巩固/训练/检测"的环节分钟数；
    3) 兜底 = 课时/3（夹在 5–20′）。
    旧实现档B 硬编码"时间 15′"、档A 显示整节课时长，两者都不对且互相打架。"""
    v = L.get('练习时长')
    if v:
        try:
            return int(v)
        except Exception:
            pass
    for ph in (L.get('环节') or []):
        if any(k in str(ph.get('名', '')) for k in ('练习', '巩固', '训练', '检测')):
            try:
                return int(ph['分'])
            except Exception:
                pass
    try:
        return min(20, max(5, int(L.get('课时', 40)) // 3))
    except Exception:
        return 15


def make_exercise_docx(path):
    doc = Document(); sec = doc.sections[0]
    sec.top_margin = sec.bottom_margin = Cm(2); sec.left_margin = sec.right_margin = Cm(2.5)
    title(doc, "课 堂 分 层 练 习", size=18)
    para(doc, LESSON['课题'] + " · 附答案（老师用，答案不下发）", size=11, color=GRAY, align=WD_ALIGN_PARAGRAPH.CENTER)
    para(doc, f"班级：＿＿＿＿　姓名：＿＿＿＿　　时间 {exercise_minutes(LESSON)}′", size=11, space_after=8)
    for group, items in EXERCISES:
        rich_para(doc, [(group, True, GRAY)], size=12)
        for no, stem, _ans in items:
            para(doc, f"{no}. {stem}", size=12, space_after=22)
        doc.add_paragraph()
    doc.add_page_break()
    heading(doc, "参考答案", size=14, color=RED)
    for group, items in EXERCISES:
        rich_para(doc, [(group, True, GRAY)], size=11)
        for no, _stem, ans in items:
            para(doc, f"{no}. {ans}", size=11.5, space_after=3)
    doc.save(path); return path

def make_homework_docx(path):
    doc = Document(); sec = doc.sections[0]
    sec.top_margin = sec.bottom_margin = Cm(2); sec.left_margin = sec.right_margin = Cm(2.5)
    title(doc, "课 后 作 业", size=18)
    para(doc, LESSON['课题'] + " · " + LESSON['作业说明'], size=11, color=GRAY, align=WD_ALIGN_PARAGRAPH.CENTER)
    para(doc, "班级：＿＿＿＿　姓名：＿＿＿＿　上交时间：＿＿＿＿", size=11, space_after=8)
    for group, items in HOMEWORK.items():
        rich_para(doc, [(group, True, GRAY)], size=12)
        for no, stem, _ans in items:
            para(doc, f"{no}. {stem}", size=12, space_after=18)
        doc.add_paragraph()
    doc.add_page_break()
    heading(doc, "参考答案（老师用）", size=14, color=RED)
    for group, items in HOMEWORK.items():
        for no, _stem, ans in items:
            para(doc, f"{no}. {ans}", size=11.5, space_after=3)
    doc.save(path); return path

def _set_run(r, size, bold=False, color=None, name='楷体'):
    r.font.size = PPt(size); r.font.bold = bold
    if color is not None: r.font.color.rgb = color
    r.font.name = name
    # 中文字体
    ea = r.font._rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(r.font._rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', name)

def make_ppt(path):
    """数据驱动的通用课件：封面 / 学习目标 / 重难点 / 导入 / 各环节一页 / 课堂小结 / 作业收尾。
    页数随环节数动态变化，任意学科可用；视觉对齐 lesson-pack「真实教案」风格：
    白底 + 学科色块标题栏 + 黑墨主体 + 红笔标重点，不用花哨渐变。
    （早期版本有天平示意、四步流程、例题板书等理科专属页，通用化后已移除，
      相关残留函数 big_formula / para_tf 一并清理，避免误导后来人。）
    """
    if not PPT_OK:
        return None
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN as _PPA
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    # 主题配色（由 --theme 设定，默认数学深蓝；红笔/淡红底为约定警示色）
    _T = (CUR or apply_theme('数学'))['ppt']
    INKC   = _T['ink']
    SUB    = _T['sub']
    HEAD   = _T['head']     # 主色：标题栏/封面
    ACC    = _T['acc']      # 强调
    REDC   = _T['red']      # 红笔批注
    PALE   = _T['pale']     # 主色淡底
    PALE_R = _T['pale_r']   # 警示淡底
    PAPER  = PRGB(0xFF,0xFF,0xFF)
    WHITE  = _T['white']
    TAG    = _T['tag']

    def bg(s):
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0,SW,SH)
        r.fill.solid(); r.fill.fore_color.rgb = PAPER; r.line.fill.background()
        r.shadow.inherit = False
        s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
        return r

    def title_bar(s, txt, stage=""):
        """顶部深蓝标题栏 + 右侧环节小标签"""
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, Inches(0.95))
        bar.fill.solid(); bar.fill.fore_color.rgb = HEAD; bar.line.fill.background(); bar.shadow.inherit=False
        tb = bar.text_frame; tb.margin_left=Inches(0.5); tb.margin_top=Inches(0.06); tb.margin_bottom=Inches(0.06)
        tb.word_wrap=False; tb.vertical_anchor = MSO_ANCHOR.MIDDLE
        p=tb.paragraphs[0]; p.text=txt; p.alignment=_PPA.LEFT
        for r in p.runs: _set_run(r, 26, True, WHITE, '黑体')
        if stage:
            tag = s.shapes.add_textbox(SW-Inches(3.0), Inches(0.22), Inches(2.7), Inches(0.5))
            tf=tag.text_frame; tf.word_wrap=False; tp=tf.paragraphs[0]; tp.text=stage; tp.alignment=_PPA.RIGHT
            for r in tp.runs: _set_run(r, 13, True, TAG)

    def note_bar(s, stage, timemin):
        """页脚：归属环节+建议时长"""
        nb = s.shapes.add_textbox(Inches(0.5), SH-Inches(0.55), SW-Inches(1.0), Inches(0.4))
        nf=nb.text_frame; nf.word_wrap=False; np=nf.paragraphs[0]; np.text = f"{stage}" + (f" · 建议 {timemin}′" if timemin else "")
        for r in np.runs: _set_run(r, 11, False, SUB)

    def bullet_box(s, x, y, w, h, title, rows, size=20, title_color=HEAD, fill=PALE):
        """圆角色块容器(教学要点/概念卡片)"""
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.color.rgb=ACC; box.line.width=PPt(0.75)
        box.shadow.inherit=False
        tf = box.text_frame; tf.word_wrap=True
        tf.margin_left=Inches(0.22); tf.margin_right=Inches(0.18)
        tf.margin_top=Inches(0.12); tf.margin_bottom=Inches(0.1)
        tp=tf.paragraphs[0]; tp.text=title; tp.space_after=PPt(6)
        for r in tp.runs: _set_run(r, size+3, True, title_color)
        first=True
        for row in rows:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first=False; p.space_after=PPt(4)
            r=p.add_run(); r.text=("• " if not row.startswith(("①","②","③","④")) else "") + row
            _set_run(r, size, False, INKC)
        return box

    # ============================================================
    # 课件内容组装 —— 完全由单源 LESSON 驱动，任意学科/环节数通用。
    # 页结构：封面 / 目标 / 重难点 / 导入(环节[0]) / 各环节一页 / 作业收尾
    # ============================================================
    L = LESSON
    def _cn(n):
        CN = ["零","一","二","三","四","五","六","七","八","九"]
        n = int(n)
        if n < 10:
            return CN[n]
        if n == 10:
            return "十"
        if n < 20:
            return "十" + CN[n-10]
        # 20 及以上：旧实现 CN[n-10] 会越界崩溃。环节数不会真到 20，
        # 但宁可退化成阿拉伯数字也别让脚本挂掉。
        return str(n)
    def clip(s, n=6):
        segs = [x.strip() for x in str(s).replace('①','。①').replace('②','。②')
                .replace('③','。③').replace('④','。④').split('。') if x.strip()]
        return segs[:n]
    def rowmap(ph):
        return {k: v for k, v in ph['rows']}

    # ---- P1 封面 ----
    s = prs.slides.add_slide(blank); bg(s)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SW, Inches(2.3))
    band.fill.solid(); band.fill.fore_color.rgb=HEAD; band.line.fill.background(); band.shadow.inherit=False
    big = s.shapes.add_textbox(Inches(0.6), Inches(2.75), SW-Inches(1.2), Inches(1.15))
    btf=big.text_frame; bp=btf.paragraphs[0]; bp.text=L['课题']; bp.alignment=_PPA.CENTER
    for r in bp.runs: _set_run(r, 52, True, WHITE, '黑体')
    sub = s.shapes.add_textbox(Inches(0.6), Inches(4.05), SW-Inches(1.2), Inches(0.6))
    sp=sub.text_frame.paragraphs[0]; sp.text=L['章节']; sp.alignment=_PPA.CENTER
    for r in sp.runs: _set_run(r, 20, False, TAG)
    meta = s.shapes.add_textbox(Inches(0.6), Inches(5.35), SW-Inches(1.2), Inches(0.9))
    mp=meta.text_frame; mp.word_wrap=True
    for i,ln in enumerate([f"{L['学科']} · {L['年级']} · {L['教材']}", f"{L['课型']} · {L['课时']} 分钟"]):
        p=mp.paragraphs[0] if i==0 else mp.add_paragraph(); p.alignment=_PPA.CENTER
        for r in (p.add_run(),): r.text=ln; _set_run(r, 18, False, SUB)
    top=s.shapes.add_textbox(Inches(0.6), Inches(0.7), SW-Inches(1.2), Inches(0.6))
    tpp=top.text_frame.paragraphs[0]; tpp.text=L['章节']; tpp.alignment=_PPA.CENTER
    for r in tpp.runs: _set_run(r, 20, True, HEAD, '宋体')

    # ---- P2 学习目标 ----
    s=prs.slides.add_slide(blank); bg(s); title_bar(s,"本课学习目标","目标")
    goals = []
    for cat, items in L['教学目标']:
        for j, it in enumerate(items):
            goals.append((cat if j==0 else "", it))
    goals = goals[:6]
    y0=1.5; step=0.95; hh=0.9
    for i,(cat,it) in enumerate(goals):
        bb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y0+i*step), SW-Inches(2.0), Inches(hh))
        bb.fill.solid(); bb.fill.fore_color.rgb=PALE; bb.line.color.rgb=ACC; bb.line.width=PPt(0.75); bb.shadow.inherit=False
        tf=bb.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.3); tf.margin_top=Inches(0.05)
        p=tf.paragraphs[0]; p.text=f"{i+1}　{'【'+cat+'】' if cat else ''}{it}"
        for r in p.runs: _set_run(r, 17, True, HEAD)
    note_bar(s, "目标", None)

    # ---- P3 教学重难点 ----
    kd = L['重难点']
    s=prs.slides.add_slide(blank); bg(s); title_bar(s,"教学重难点","本课关键")
    bullet_box(s, Inches(1.0), Inches(1.7), SW-Inches(2.0), Inches(1.6),
        "重点", clip(kd[0],3), size=19)
    bullet_box(s, Inches(1.0), Inches(3.45), SW-Inches(2.0), Inches(1.6),
        "难点", clip(kd[1],3), size=19)
    bullet_box(s, Inches(1.0), Inches(5.2), SW-Inches(2.0), Inches(1.7),
        "难点 · 突破策略", clip(kd[2],3), size=17, fill=PALE_R, title_color=REDC)

    # ---- 通用环节页 ----
    def phase_page(ph, stage_label):
        d = rowmap(ph)
        s=prs.slides.add_slide(blank); bg(s); title_bar(s, ph['名'], stage_label)
        teach = clip(d.get('教师活动',''), 6)
        right = []
        sv = d.get('学生活动'); cv = d.get('当堂检验')
        if sv: right += ["【学生活动】"+s2 for s2 in clip(sv,2)]
        if cv: right += ["【当堂检验】"+s2 for s2 in clip(cv,2)]
        if not right: right = ["（见教学设计）"]
        bullet_box(s, Inches(1.0), Inches(1.5), Inches(6.4), Inches(5.2),
            "教师活动", teach, size=19)
        bullet_box(s, Inches(7.7), Inches(1.5), Inches(4.8), Inches(5.2),
            "学生活动 · 当堂检验", right, size=16, fill=PALE)
        note_bar(s, stage_label, ph['分'])
        return s

    # ---- P4 导入（环节[0]）----
    n_ph = len(L['环节'])
    phase_page(L['环节'][0], f"导入 · {L['环节'][0]['分']}′")

    # ---- 中间教学内容环节（索引1起，到倒数第二为止；小结与作业各自成专题页）----
    for idx, ph in enumerate(L['环节'][1:-1], start=2):
        phase_page(ph, f"环节{_cn(idx)} · {ph['分']}′")

    # ---- 课堂小结总览 ----
    # 旧实现固定取"倒数第二个环节"，当小结与作业合并成最后一个环节时（如语文《咏雪》的
    # "五、小结与作业布置"），倒数第二个其实是"品读新授"，小结页会显示成新授内容。
    # 改为按名称定位小结环节，找不到再回退到倒数第二个。
    s=prs.slides.add_slide(blank); bg(s); title_bar(s,"课堂小结 · 一图回顾","小结")
    recap_src = None
    for ph in reversed(L['环节']):
        if any(k in str(ph.get('名', '')) for k in ('小结', '总结', '回顾')):
            recap_src = ph
            break
    if recap_src is None and n_ph >= 2:
        recap_src = L['环节'][-2]
    recap = clip(rowmap(recap_src).get('教师活动', ''), 4) if recap_src else ["（见教学设计）"]
    bullet_box(s, Inches(1.0), Inches(1.8), Inches(6.4), Inches(3.4),
        "一起回顾", recap, size=18)
    bullet_box(s, Inches(7.6), Inches(1.8), Inches(5.0), Inches(3.4),
        "本课核心（重难点）", clip(kd[0],3), size=18, fill=PALE)
    board_txt = L.get('板书', {}).get('title', L['课题'])
    rb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.3), SW-Inches(2.0), Inches(1.3))
    rb.fill.solid(); rb.fill.fore_color.rgb=PALE_R; rb.line.color.rgb=REDC; rb.line.width=PPt(1.2); rb.shadow.inherit=False
    tf=rb.text_frame; tf.word_wrap=True; tf.margin_top=Inches(0.1)
    p=tf.paragraphs[0]; p.text="板书主题：" + str(board_txt)
    for r in p.runs: _set_run(r, 20, True, REDC, '黑体')

    # ---- 作业收尾（环节[-1]，通常是作业布置）----
    # 旧实现读 L.get('homework')：homework 是顶层键、不在 lesson 里，永远取不到，
    # 导致作业页始终退化成"教师活动"文本，列出的作业组/题数从未生效过。
    # 只有 1 个环节时，环节[0] 已在导入页出过，再出一页收尾会重复同一环节
    if n_ph >= 2:
        lastph = L['环节'][-1]
        d = rowmap(lastph)
        s=prs.slides.add_slide(blank); bg(s); title_bar(s, lastph['名'], "收尾")
        hw_rows = []
        if isinstance(HOMEWORK, dict) and HOMEWORK:
            for g, items in HOMEWORK.items():
                hw_rows.append(f"《{g}》 · {len(items)} 题")
        if not hw_rows:
            hw_rows = clip(d.get('教师活动',''), 4)
        bullet_box(s, Inches(1.0), Inches(1.8), SW-Inches(2.0), Inches(3.4),
            "课后作业（详见作业纸）", hw_rows, size=20)
        tips = clip(d.get('教师活动',''), 2)
        nx=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.4), SW-Inches(2.0), Inches(1.3))
        nx.fill.solid(); nx.fill.fore_color.rgb=PALE; nx.line.color.rgb=ACC; nx.line.width=PPt(1); nx.shadow.inherit=False
        tf=nx.text_frame; tf.margin_top=Inches(0.12); tf.word_wrap=True
        p=tf.paragraphs[0]; p.text="收尾：" + (tips[0] if tips else "见教学设计")
        for r in p.runs: _set_run(r, 18, True, ACC)
        note_bar(s, "收尾", lastph['分'])

    prs.save(path); return path

def main():
    ap = argparse.ArgumentParser(description="生成整套课时教学包（教案/课件/学案/练习/作业）")
    ap.add_argument('--out', default='.', help='输出目录')
    ap.add_argument('--theme', default='数学',
                    help="主色主题：预设名(数学/语文/英语/物理/化学/生物/历史/地理/通用) "
                         "或 '#hex' 自定义主色，例如 --theme 语文 或 --theme '#C0392B'")
    ap.add_argument('--data', default=None,
                    help="单源内容 JSON 路径（结构见 templates/lesson-content.schema.md）。"
                         "缺省用内置数学《一元一次方程》演示样例。")
    ap.add_argument('--strict-time', action='store_true',
                    help="时间账（环节分钟合计≠课时）不平时直接终止，不生成任何文件。"
                         "默认仅告警后继续。")
    a = ap.parse_args()
    # 外部单源 JSON → 替换内置数据结构，实现任意学科/课题通用化
    if a.data:
        load_data(a.data)
    # 先应用主题：整套 docx + pptx 的主色统一走主题，红笔警示色保持约定
    apply_theme(a.theme)
    os.makedirs(a.out, exist_ok=True)
    g = os.path.join

    # 时间账硬校验（quality-checklist A 项）：不平要显式说出来，不能静默产出
    ok, issues = check_time(LESSON)
    if not ok:
        print("  ⚠ 时间账不平（quality-checklist A 项，请修订内容 JSON）：")
        for it in issues:
            print("      - " + it)
        if a.strict_time:
            raise SystemExit("✗ 已启用 --strict-time，时间账不平即终止，未生成任何文件。")

    # 产物 → 依赖的顶层数据键。键缺失/为空则跳过该产物并说明，
    # 绝不退回内置数学样例（否则会出现"语文课配方程作业"的串味事故）。
    files = [
        ("教学设计.docx",   make_design_docx,   'lesson',          bool(LESSON)),
        ("学习任务单.docx", make_worksheet_docx,'worksheet_tasks', bool(WORKSHEET_TASKS)),
        ("课堂分层练习.docx", make_exercise_docx, 'exercises',     bool(EXERCISES)),
        ("课后作业.docx",   make_homework_docx, 'homework',        bool(HOMEWORK)),
        ("教学课件.pptx",   make_ppt,           'lesson',          bool(LESSON)),
    ]
    made, skipped, failed = [], [], []
    for name, fn, key, ready in files:
        if not ready:
            skipped.append(f"{name}（内容 JSON 未提供 '{key}'）")
            continue
        try:
            p = fn(g(a.out, name))
            if p:
                made.append(p)
            else:
                skipped.append(f"{name}（缺 python-pptx，已优雅跳过）")
        except Exception as e:
            failed.append(name)
            print(f"  ✗ {name} 生成失败：{e}")

    src = ('外部 ' + a.data) if a.data else '内置数学《一元一次方程》演示样例'
    print(f"主题：{CUR['name']}；数据源：{src}；已生成 {len(made)} 个文件：")
    for p in made:
        print("  " + p)
    if skipped:
        print("  已跳过：")
        for s in skipped:
            print("    - " + s)
    if not a.data:
        print("  提示：未传 --data，当前产出为内置演示样例，请勿直接用于教学。")

    # 失败必须以非 0 退出码暴露，避免自动化/脚本误判为"全部成功"而拿到残缺包
    if failed:
        print(f"✗ {len(failed)} 个产物生成失败：{', '.join(failed)}")
        sys.exit(1)

if __name__ == '__main__':
    main()
